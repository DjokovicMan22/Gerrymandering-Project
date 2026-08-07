#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

import pandas as pd

from redistricting.reporting.week7 import (
    clean_metric_name,
    extract_percentile_rows,
    format_number,
    get_metric_row,
    read_csv_if_exists,
    top_sorted_deviations,
)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a stronger evidence-driven Week 7 presentation deck.")
    parser.add_argument("--name", default="mi_2020")
    parser.add_argument("--state-name", default="Michigan")
    parser.add_argument("--election-label", default="2020 presidential election")
    return parser.parse_args()


def add_title(slide, text, x, y, w, h, size=30, color=None):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    if color:
        p.font.color.rgb = RGBColor.from_string(color)
    return tb


def add_text(slide, text, x, y, w, h, size=16, bold=False, color=None):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = RGBColor.from_string(color)
    return tb


def add_bullets(slide, items, x, y, w, h, size=15):
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
    return tb


def add_image_or_note(slide, path, x, y, w, h=None, note=None):
    from pptx.util import Inches
    p = Path(path)
    if p.exists():
        if h is None:
            slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w))
        else:
            slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        add_text(slide, note or f"Missing figure: {p}", x, y, w, h or 1.0, size=14, bold=True, color="8A1C1C")


def add_footer(slide, n):
    add_text(slide, f"{n}/16", 12.1, 7.05, 0.55, 0.25, size=8, color="666666")


def metric_lines(percentiles):
    rows = extract_percentile_rows(percentiles)
    if not rows:
        return ["Percentile table not parsed"]
    out = []
    for r in rows:
        out.append(f"{clean_metric_name(r.metric)}: enacted {format_number(r.enacted_value, 3)}, percentile {format_number(r.percentile, 1)}%")
    return out


def main() -> None:
    args = parse_args()
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
    except Exception as exc:
        raise SystemExit("Install python-pptx first: python -m pip install python-pptx") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    percentiles = read_csv_if_exists(f"outputs/tables/week5/{args.name}_outlier_percentiles.csv")
    diagnostics = read_csv_if_exists(f"outputs/tables/week4/{args.name}_production_diagnostics.csv")
    sorted_summary = read_csv_if_exists(f"outputs/tables/week5/{args.name}_sorted_district_summary.csv")
    geography = read_csv_if_exists(f"outputs/tables/week6/{args.name}_geography_baseline.csv")
    ranges = read_csv_if_exists(f"outputs/tables/week6/{args.name}_percentile_ranges.csv")
    deviations = top_sorted_deviations(sorted_summary, 3)

    def new_slide(n, title, subtitle=None):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor.from_string("F7F5EF")
        add_title(slide, title, 0.55, 0.25, 11.8, 0.45, size=27)
        if subtitle:
            add_text(slide, subtitle, 0.58, 0.78, 10.8, 0.28, size=12, color="555555")
        add_footer(slide, n)
        return slide

    # 1
    s = new_slide(1, f"Detecting Outlier Maps with Ensembles", f"{args.state_name} congressional plan • {args.election_label}")
    add_text(s, "Research question", 0.75, 1.55, 3.0, 0.35, size=18, bold=True)
    add_text(s, f"Where does the enacted {args.state_name} plan fall relative to thousands of valid alternative maps sampled under explicit neutral constraints?", 0.75, 2.02, 6.4, 1.15, size=25, bold=True)
    add_bullets(s, ["Not a visual intuition project", "Not a legal verdict", "A reproducible statistical baseline"], 0.95, 4.0, 5.3, 1.2, size=17)
    add_image_or_note(s, f"outputs/figures/week5/{args.name}_signature_outlier_panel.png", 7.35, 1.25, 5.25, 4.35)

    # 2
    s = new_slide(2, "Hypothesis and claim standard", "A strong claim needs converging evidence, not one metric.")
    add_text(s, "Hypothesis", 0.75, 1.35, 2.0, 0.3, size=17, bold=True)
    add_text(s, "If the enacted plan is structurally unusual, it should fall in the tails of the sampled ensemble across multiple forms of evidence: seats, efficiency gap, mean-median, sorted districts, and robustness tests.", 0.75, 1.8, 11.6, 1.2, size=22, bold=True)
    add_bullets(s, ["Evidence required: percentile tails", "Structure required: sorted-district deviations", "Credibility required: diagnostics + sensitivity", "Boundary required: explicit limitations"], 1.0, 3.55, 11.2, 1.8, size=18)

    # 3
    s = new_slide(3, "Research context: where this project fits")
    add_text(s, "The paper is stronger when it is explicitly connected to existing computational-redistricting literature.", 0.8, 1.15, 11.5, 0.55, size=23, bold=True)
    add_bullets(s, [
        "Guest et al. (2019): computational redistricting makes criteria explicit; compactness is only one possible criterion.",
        "Stephanopoulos (2017): efficiency gap captures packing/cracking but causes and consequences require institutional context.",
        "Kenny et al. (2023): compare enacted plans to simulated nonpartisan baselines to separate map effects from geography.",
        "Bouton et al. (2024): turnout heterogeneity creates pack-crack-pack strategies ordinary metrics may miss."
    ], 0.95, 2.05, 11.4, 2.8, size=18)
    add_text(s, "Positioning: this project is a state-level, reproducible ensemble baseline—not a legal verdict, national causal estimate, or pure compactness optimization.", 1.0, 5.55, 11.2, 0.7, size=19, bold=True)

    # 4
    s = new_slide(4, "Why visual compactness is not enough")
    add_bullets(s, ["Irregular districts can come from water, boundaries, or legal constraints.", "Compact districts can still produce partisan advantage.", "Compactness metrics are useful diagnostics, not fairness definitions.", "This motivates comparing enacted maps against an ensemble baseline."], 0.8, 1.3, 5.7, 4.6, size=18)
    add_image_or_note(s, "outputs/figures/week2/boundary_resolution_sensitivity.png", 6.9, 1.25, 5.7, 4.4)

    # 4
    s = new_slide(5, "Data model: a map becomes a graph")
    add_bullets(s, ["Nodes = precincts / voting tabulation districts", "Edges = shared borders using rook adjacency", "Node attributes = population, enacted district, election votes", "Disconnected island components repaired with documented artificial bridge edges"], 0.75, 1.25, 5.8, 4.5, size=18)
    add_image_or_note(s, "outputs/figures/week1/mi_dual_graph_overlay.png", 6.8, 1.15, 5.7, 4.7, note="Use your Week 1 dual graph map here.")

    # 5
    s = new_slide(6, "ReCom ensemble design")
    add_bullets(s, ["Start at enacted plan (step 0)", "Merge two adjacent districts", "Re-split with a random spanning-tree procedure", "Accept plans satisfying population and contiguity constraints", "Record seat count and partisan metrics at every step"], 0.75, 1.25, 5.6, 4.8, size=18)
    add_text(s, "Main run", 7.05, 1.35, 2.2, 0.3, size=17, bold=True)
    add_bullets(s, ["3 independent chains", "7,000 steps each", "700-step burn-in removed", "18,900 post-burn-in plans", "2% population tolerance baseline"], 7.1, 1.85, 4.9, 2.7, size=22)

    # 6
    s = new_slide(7, "Diagnostics before interpretation", "A completed run is not automatically a credible ensemble.")
    add_image_or_note(s, f"outputs/figures/week4/{args.name}_production_dem_seats.png", 0.7, 1.1, 5.8, 4.8)
    add_image_or_note(s, f"outputs/figures/week4/{args.name}_production_efficiency_gap.png", 6.85, 1.1, 5.8, 4.8)
    if diagnostics is not None and "split_rhat" in diagnostics.columns:
        add_text(s, f"Max split-Rhat: {format_number(diagnostics['split_rhat'].max(), 3)}", 0.9, 6.35, 4.0, 0.35, size=15, bold=True)
    add_text(s, "Interpretation: reassuring practical diagnostics, not proof of perfect mixing.", 5.1, 6.35, 6.7, 0.35, size=15, bold=True)

    # 7
    s = new_slide(8, "Main evidence: empirical percentiles")
    add_image_or_note(s, f"outputs/figures/week5/{args.name}_signature_outlier_panel.png", 0.65, 1.05, 7.2, 5.4)
    add_bullets(s, metric_lines(percentiles), 8.15, 1.35, 4.6, 3.8, size=16)
    add_text(s, "Percentile = enacted plan's position inside the sampled post-burn-in ensemble, not a classical p-value.", 8.15, 5.4, 4.45, 0.8, size=14, bold=True)

    # 8
    s = new_slide(9, "Sorted-district plot: structure behind the aggregate")
    add_image_or_note(s, f"outputs/figures/week5/{args.name}_sorted_districts.png", 0.65, 1.0, 8.0, 5.45)
    if deviations is not None:
        lines = []
        for _, r in deviations.iterrows():
            lines.append(f"Rank {int(r['rank'])}: {format_number(100*r['enacted_minus_ensemble_median'], 2, signed=True)} pp vs ensemble median")
        add_bullets(s, lines, 9.0, 1.55, 3.6, 2.0, size=15)
    add_text(s, "Why it matters: packing and cracking are visible as rank-level deviations, not just total seats.", 9.0, 4.5, 3.6, 1.1, size=15, bold=True)

    # 9
    s = new_slide(10, "Maps: statistics tied back to geography")
    add_image_or_note(s, f"outputs/figures/week5/{args.name}_plan_comparison_maps.png", 0.55, 1.0, 12.2, 5.75)

    # 10
    s = new_slide(11, "Evidence ledger", "What each result supports — and what it does not prove.")
    rows = extract_percentile_rows(percentiles)[:4]
    y = 1.3
    for r in rows:
        add_text(s, clean_metric_name(r.metric), 0.75, y, 3.0, 0.25, size=15, bold=True)
        add_text(s, f"enacted {format_number(r.enacted_value, 3)} • percentile {format_number(r.percentile, 1)}%", 3.85, y, 4.1, 0.25, size=15)
        add_text(s, "supports conditional outlier analysis", 8.1, y, 4.2, 0.25, size=15)
        y += 0.7
    add_text(s, "None of these alone proves intent or illegality. The strength comes from cumulative evidence plus limitations.", 0.9, 5.45, 11.3, 0.55, size=19, bold=True)

    # 11
    s = new_slide(12, "Stress test 1: population tolerance")
    add_image_or_note(s, f"outputs/figures/week6/{args.name}_constraint_distributions.png", 0.65, 1.0, 6.1, 5.2)
    add_image_or_note(s, f"outputs/figures/week6/{args.name}_constraint_percentiles.png", 7.0, 1.0, 5.65, 5.2)
    if ranges is not None and "percentile_range" in ranges.columns:
        add_text(s, f"Largest percentile movement: {format_number(ranges['percentile_range'].max(), 2)} points", 0.9, 6.45, 5.5, 0.3, size=15, bold=True)

    # 12
    s = new_slide(13, "Stress test 2: geography baseline")
    add_image_or_note(s, f"outputs/figures/week6/{args.name}_geography_baseline.png", 0.7, 1.1, 7.1, 5.2)
    if geography is not None and not geography.empty:
        g = geography.iloc[0]
        bullets = []
        for col, label in [("proportional_dem_seats", "Proportional benchmark"), ("ensemble_mean_dem_seats", "Ensemble mean"), ("enacted_dem_seats", "Enacted")]:
            if col in geography.columns:
                bullets.append(f"{label}: {format_number(g[col], 2)} Democratic seats")
        add_bullets(s, bullets, 8.25, 1.7, 4.1, 2.2, size=19)
    add_text(s, "This avoids the weak argument that any disproportionality equals gerrymandering.", 8.25, 4.6, 4.1, 1.0, size=16, bold=True)

    # 13
    s = new_slide(14, "What the paper can and cannot claim")
    add_text(s, "Supported", 0.9, 1.25, 2.0, 0.3, size=18, bold=True)
    add_bullets(s, ["Empirical outlier status relative to sampled ensemble", "Transparent assumptions and constraints", "Practical multi-chain diagnostics", "Robustness checks across tolerances"], 0.9, 1.75, 5.0, 3.0, size=17)
    add_text(s, "Not supported", 7.0, 1.25, 2.0, 0.3, size=18, bold=True)
    add_bullets(s, ["Legal conclusion", "Proof of intent", "Proof of perfect mixing", "Full modeling of every legal criterion", "A universal result across all elections"], 7.0, 1.75, 5.0, 3.4, size=17)

    # 14
    s = new_slide(15, "Strongest counterargument")
    add_text(s, "The ensemble may not represent every legally realistic plan because it does not fully encode Voting Rights Act compliance, communities of interest, county splits, incumbency, or every state-specific rule.", 0.9, 1.45, 11.4, 1.35, size=25, bold=True)
    add_bullets(s, ["Response: disclose constraints instead of hiding them", "Run sensitivity tests", "Avoid claims about intent or legality", "Treat this as measurement evidence, not a verdict"], 1.0, 3.45, 10.5, 2.0, size=19)

    # 15
    s = new_slide(16, "Final conclusion")
    add_text(s, f"Under the stated graph, data, ReCom proposal, burn-in, population constraints, and {args.election_label} scoring data, the enacted {args.state_name} plan falls at the reported empirical percentiles of the sampled ensemble.", 0.9, 1.25, 11.4, 1.5, size=25, bold=True)
    add_bullets(s, ["This is stronger than visual intuition.", "It is more honest than a single metric.", "It becomes persuasive because assumptions, diagnostics, and limitations are explicit."], 1.0, 3.55, 10.8, 1.6, size=21)

    out = Path(f"outputs/slides/week7/{args.name}_talk.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Saved {out}")
    print("Deck is 15 slides and evidence-driven. Inspect visually and trim for a 10-minute talk if needed.")


if __name__ == "__main__":
    main()
